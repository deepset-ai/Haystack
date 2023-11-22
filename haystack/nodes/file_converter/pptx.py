from haystack.nodes.base import BaseComponent
from haystack.schema import Document
from haystack.lazy_imports import LazyImport
from haystack.nodes.file_converter.base import BaseConverter
from typing import Tuple, List, Optional, Any, Dict
from pathlib import Path
import logging


logger = logging.getLogger(__name__)


with LazyImport("Run 'pip install python-pptx'") as pptx_import:
    from pptx import Presentation


class PptxConverter(BaseConverter):
  def __init__(
        self,
        remove_numeric_tables: bool = False,
        valid_languages: Optional[List[str]] = None,
        id_hash_keys: Optional[List[str]] = None,
        progress_bar: bool = True,
    ):
        pptx_import.check()
        super().__init__(
            remove_numeric_tables=remove_numeric_tables,
            valid_languages=valid_languages,
            id_hash_keys=id_hash_keys,
            progress_bar=progress_bar,
        )

  def convert(
        self,
        file_path: Path,
        meta: Optional[Dict[str, str]] = None,
        remove_numeric_tables: Optional[bool] = None,
        valid_languages: Optional[List[str]] = None,
        encoding: Optional[str] = None,
        id_hash_keys: Optional[List[str]] = None,
    ) -> List[Document]:
        """
        Extract text from a .pptx file.
        Note: As pptx doesn't contain "page" information, we actually extract and return a list of texts from each slide here.
        For compliance with other converters we nevertheless opted for keeping the methods name.

        :param file_path: Path to the .pptx file you want to convert
        :param meta: dictionary of meta data key-value pairs to append in the returned document.
        :param remove_numeric_tables: This option uses heuristics to remove numeric rows from the tables.
                                      The tabular structures in documents might be noise for the reader model if it
                                      does not have table parsing capability for finding answers. However, tables
                                      may also have long strings that could possible candidate for searching answers.
                                      The rows containing strings are thus retained in this option.
        :param valid_languages: validate languages from a list of languages specified in the ISO 639-1
                                (https://en.wikipedia.org/wiki/ISO_639-1) format.
                                This option can be used to add test for encoding errors. If the extracted text is
                                not one of the valid languages, then it might likely be encoding error resulting
                                in garbled text.
        :param encoding: Not applicable
        :param id_hash_keys: Generate the document id from a custom list of strings that refer to the document's
            attributes. If you want to ensure you don't have duplicate documents in your DocumentStore but texts are
            not unique, you can modify the metadata and pass e.g. `"meta"` to this field (e.g. [`"content"`, `"meta"`]).
            In this case the id will be generated by using the content and the defined metadata.
        """
        if remove_numeric_tables is None:
            remove_numeric_tables = self.remove_numeric_tables
        if valid_languages is None:
            valid_languages = self.valid_languages
        if remove_numeric_tables is True:
            raise Exception("'remove_numeric_tables' is not supported by PptxToTextConverter.")
        if valid_languages is True:
            raise Exception("Language validation using 'valid_languages' is not supported by PptxToTextConverter.")
        if id_hash_keys is None:
            id_hash_keys = self.id_hash_keys
          
        pres = Presentation(file_path)
        text = ""
        for slide in pres.slides:
          for shape in slide.shapes:
            if hasattr(shape, "text"):
              text += shape.text
              
        document = Document(content=text, meta=meta, id_hash_keys=id_hash_keys)
        return [document]  
